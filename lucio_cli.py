import click
import os
import time
import sys
import logging
import concurrent.futures
from pathlib import Path
from typing import List, Tuple, Dict, Optional, Set
import gc
import json
import re
import math
import numpy as np
from collections import Counter

os.environ["OMP_NUM_THREADS"] = "1"
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

MAX_WORKERS = 8
DATA_DIR = "./data"
CHUNKS_PATH = "chunks.json"
META_PATH = "metadata.json"
TFIDF_PATH = "tfidf_index.json"
TFIDF_VECTORS_PATH = "tfidf_vectors.npz"
VOCAB_PATH = "tfidf_vocab.json"
RESULTS_PATH = "results.json"

START_TIME = time.time()

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "AIzaSyA9z88EYfTuKuG78TtBIOzJIiM6z825dM8")
GEMINI_MODEL = "gemini-2.5-flash"
GEMINI_URL = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent"

CHUNK_SIZE = 400
CHUNK_OVERLAP = 80
MAX_CHUNKS_PER_DOC = 8
MAX_PDF_PAGES = 80


def adaptive_optimization(doc_count: int):
    global CHUNK_SIZE, CHUNK_OVERLAP, MAX_CHUNKS_PER_DOC, MAX_PDF_PAGES
    if doc_count <= 20:
        CHUNK_SIZE, CHUNK_OVERLAP, MAX_CHUNKS_PER_DOC, MAX_PDF_PAGES = 500, 100, 30, 80
    elif doc_count <= 50:
        CHUNK_SIZE, CHUNK_OVERLAP, MAX_CHUNKS_PER_DOC, MAX_PDF_PAGES = 450, 90, 20, 80
    elif doc_count <= 100:
        CHUNK_SIZE, CHUNK_OVERLAP, MAX_CHUNKS_PER_DOC, MAX_PDF_PAGES = 400, 70, 12, 80
    elif doc_count <= 200:
        CHUNK_SIZE, CHUNK_OVERLAP, MAX_CHUNKS_PER_DOC, MAX_PDF_PAGES = 350, 60, 8, 80
    else:
        CHUNK_SIZE, CHUNK_OVERLAP, MAX_CHUNKS_PER_DOC, MAX_PDF_PAGES = 300, 50, 6, 60
    logger.info(f"Adaptive: {doc_count} docs -> chunk={CHUNK_SIZE}, overlap={CHUNK_OVERLAP}, "
                f"max_chunks/doc={MAX_CHUNKS_PER_DOC}, max_pdf_pages={MAX_PDF_PAGES}")


# ==============================================================================
# GEMINI API
# ==============================================================================

def call_gemini(prompt: str, timeout: float = 20.0) -> str:
    import urllib.request, urllib.error
    url = f"{GEMINI_URL}?key={GEMINI_API_KEY}"
    payload = json.dumps({
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": 0.0, "maxOutputTokens": 1024}
    }).encode('utf-8')
    req = urllib.request.Request(url, data=payload,
                                 headers={"Content-Type": "application/json"}, method="POST")
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            data = json.loads(resp.read().decode('utf-8'))
            cands = data.get("candidates", [])
            if not cands: return ""
            parts = cands[0].get("content", {}).get("parts", [])
            if not parts: return ""
            return parts[0].get("text", "").strip()
    except urllib.error.HTTPError as e:
        try: body = e.read().decode('utf-8')[:500]
        except: body = ""
        logger.error(f"Gemini HTTP {e.code}: {body}")
    except Exception as e:
        logger.error(f"Gemini: {type(e).__name__}: {e}")
    return ""


def test_gemini() -> bool:
    r = call_gemini("Reply with: OK", timeout=8.0)
    ok = bool(r)
    logger.info(f"Gemini test: {'OK' if ok else 'FAILED'}")
    return ok


def build_qa_prompt(question: str, context_chunks: List[Dict], corpus_summary: str = "") -> str:
    """Build QA prompt with chunk metadata (doc_name, pages) included in context."""
    parts = []
    for chunk in context_chunks:
        doc_name = chunk.get("doc_name", "unknown")
        pages = chunk.get("pages", [])
        page_str = ", ".join(f"p{p}" for p in pages) if pages else "unknown"
        header = f"[Source: {doc_name} | Pages: {page_str}]"
        parts.append(f"{header}\n{chunk['text']}")
    context = "\n\n---\n\n".join(parts)
    if len(context) > 20000:
        context = context[:20000]

    corpus_section = ""
    if corpus_summary:
        corpus_section = f"""
Corpus Information (all documents in the collection):
{corpus_summary}

"""

    return f"""Answer the question using ONLY the context below.

Rules:
- Be precise. Use EXACT numbers, names, dates, terms from the context.
- Give a complete but concise answer.
- If the question asks for multiple items, list ALL of them.
- Do NOT round numbers. If context says "$39.1 billion", say "$39.1 billion".
- Do NOT add disclaimers or caveats.
- Read through ALL the context carefully before answering.
- Only say "Information not found in provided documents" if you truly cannot find any relevant information.

After your answer, on a NEW line, output the source information in this EXACT format:
SOURCE: <document_filename>, <page_number(s)>
If multiple sources, put each on its own SOURCE: line.
Example:
SOURCE: contract.pdf, p4
SOURCE: report.pdf, p12, p13

{corpus_section}Context:
{context}

Question: {question}

Answer:"""


def parse_gemini_response(response: str) -> Dict:
    """Parse Gemini response to extract answer, document names, and page numbers."""
    if not response:
        return {"answer": "", "doc_names": [], "page_numbers": []}

    lines = response.strip().split("\n")
    answer_lines = []
    source_lines = []
    for line in lines:
        if line.strip().startswith("SOURCE:"):
            source_lines.append(line.strip())
        else:
            answer_lines.append(line)

    answer = "\n".join(answer_lines).strip()
    # Remove trailing empty lines before SOURCE
    answer = answer.rstrip()

    doc_names = []
    page_numbers = []
    for sl in source_lines:
        # Parse "SOURCE: filename.pdf, p4, p5"
        content = sl[len("SOURCE:"):].strip()
        parts = [p.strip() for p in content.split(",")]
        if parts:
            doc_name = parts[0]
            if doc_name and doc_name.lower() not in ("unknown", "n/a", "none"):
                if doc_name not in doc_names:
                    doc_names.append(doc_name)
            for p in parts[1:]:
                m = re.match(r'p(\d+)', p.strip())
                if m:
                    page_numbers.append(int(m.group(1)))

    return {"answer": answer, "doc_names": doc_names, "page_numbers": sorted(set(page_numbers))}


# ==============================================================================
# SCORING — lenient fuzzy matching
# ==============================================================================

def normalize_for_comparison(text: str) -> str:
    if not text: return ""
    t = str(text).strip().lower()
    for p in ["according to the document,", "based on the context,",
              "the answer is", "answer:", "based on the provided context,",
              "based on the provided documents,"]:
        if t.startswith(p): t = t[len(p):].strip()
    t = t.rstrip('.')
    t = re.sub(r'\s+', ' ', t)
    return t.strip()


def extract_numbers(text: str) -> Set[str]:
    """Extract all numbers from text, normalized."""
    nums = set()
    for m in re.finditer(r'[\$₹]?\s*[\d,]+\.?\d*\s*(?:billion|million|thousand|cr|crore|lakh|%|b|m|k)?', text.lower()):
        n = re.sub(r'[,\s]', '', m.group()).strip()
        if n: nums.add(n)
    for m in re.finditer(r'\b\d[\d,.]*\d\b|\b\d+\b', text):
        nums.add(m.group().replace(',', ''))
    return nums


def extract_key_terms(text: str) -> Set[str]:
    """Extract proper nouns and key terms."""
    terms = set()
    for m in re.finditer(r'[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*', text):
        terms.add(m.group().lower())
    for m in re.finditer(r'\b[A-Z]{2,}\b', text):
        terms.add(m.group().lower())
    return terms


def score_answer(predicted: str, ground_truth: str) -> dict:
    pred_raw = str(predicted).strip()
    truth_raw = str(ground_truth).strip()
    pred = normalize_for_comparison(pred_raw)
    truth = normalize_for_comparison(truth_raw)

    if not pred or not truth:
        return {"exact_match": False, "contains_truth": False,
                "token_f1": 0.0, "number_match": 0.0, "key_term_match": 0.0, "composite": 0.0}

    exact = pred == truth
    pred_has_truth = truth in pred
    truth_has_pred = pred in truth

    pred_tokens = set(pred.split())
    truth_tokens = set(truth.split())
    common = pred_tokens & truth_tokens
    if common:
        prec = len(common) / len(pred_tokens)
        rec = len(common) / len(truth_tokens)
        f1 = 2 * prec * rec / (prec + rec)
    else:
        f1 = 0.0

    pred_nums = extract_numbers(pred_raw)
    truth_nums = extract_numbers(truth_raw)
    if truth_nums:
        num_hits = len(pred_nums & truth_nums)
        num_score = num_hits / len(truth_nums)
    else:
        num_score = 0.0

    pred_terms = extract_key_terms(pred_raw)
    truth_terms = extract_key_terms(truth_raw)
    if truth_terms:
        term_hits = len(pred_terms & truth_terms)
        term_score = term_hits / len(truth_terms)
    else:
        term_score = 0.0

    if exact:
        composite = 1.0
    elif pred_has_truth:
        composite = 0.9 + 0.1 * f1
    elif truth_has_pred and len(pred) > 5:
        composite = 0.75 + 0.1 * f1
    else:
        composite = max(
            f1,
            0.35 * f1 + 0.35 * num_score + 0.30 * term_score,
        )

    return {
        "exact_match": exact,
        "contains_truth": pred_has_truth,
        "token_f1": round(f1, 3),
        "number_match": round(num_score, 3),
        "key_term_match": round(term_score, 3),
        "composite": round(composite, 3),
    }


# ==============================================================================
# TEXT EXTRACTION — extract EVERYTHING possible
# ==============================================================================

def _extract_text(fp: Path) -> str:
    s = fp.suffix.lower()
    if s == '.pdf': return _extract_pdf(fp)
    elif s == '.docx': return _extract_docx(fp)
    elif s in {'.xlsx', '.xls'}: return _extract_spreadsheet(fp)
    elif s in {'.pptx', '.ppt'}: return _extract_pptx(fp)
    elif s in {'.zip', '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'}: return ""
    else:
        try:
            with open(fp, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(500_000)
        except: return ""


def _extract_pdf(fp: Path) -> str:
    try: import fitz
    except ImportError: return ""
    try:
        pages = []
        with fitz.open(fp) as doc:
            for i in range(min(len(doc), MAX_PDF_PAGES)):
                t = doc[i].get_text("text", sort=True)
                if t and t.strip():
                    pages.append(f"[Page {i+1}]\n{t}")
        return "\n\n".join(pages)
    except Exception as e:
        logger.error(f"PDF err {fp.name}: {e}"); return ""


def _extract_docx(fp: Path) -> str:
    try:
        import docx
        doc = docx.Document(fp)
        parts = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if t: parts.append(t)
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                r = " | ".join(c for c in cells if c)
                if r: rows.append(r)
            if rows: parts.append("[TABLE]\n" + "\n".join(rows) + "\n[/TABLE]")
        return "\n".join(parts)
    except Exception as e:
        logger.error(f"DOCX err {fp.name}: {e}"); return ""


def _extract_spreadsheet(fp: Path) -> str:
    text_parts = []
    try:
        import openpyxl
        wb = openpyxl.load_workbook(fp, read_only=True, data_only=True)
        for sn in wb.sheetnames[:10]:
            ws = wb[sn]
            rows = []
            for row in ws.iter_rows(max_row=1000, values_only=True):
                cells = [str(c).strip() if c is not None else "" for c in row]
                if any(c for c in cells):
                    rows.append(" | ".join(cells))
            if rows:
                text_parts.append(f"[Sheet: {sn}]\n" + "\n".join(rows))
        wb.close()
        if text_parts:
            return "\n\n".join(text_parts)
    except Exception:
        pass
    try:
        import pandas as pd
        dfs = pd.read_excel(fp, sheet_name=None, nrows=1000)
        for name, df in dfs.items():
            text_parts.append(f"[Sheet: {name}]\n{df.to_string(index=False)}")
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"XLS err {fp.name}: {e}"); return ""


def _extract_pptx(fp: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(fp)
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t: texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        texts.append(" | ".join(cells))
            if texts: parts.append(f"[Slide {i+1}] " + " ".join(texts))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"PPTX err {fp.name}: {e}"); return ""


# ==============================================================================
# CLEANING — minimal, preserve content
# ==============================================================================

def clean_text(text: str) -> str:
    text = re.sub(r'\n{4,}', '\n\n\n', text)
    text = re.sub(r'[ \t]{3,}', ' ', text)
    text = re.sub(r'(?i)this page intentionally left blank', '', text)
    return text


# ==============================================================================
# PAGE-AWARE CHUNKING WITH METADATA
# ==============================================================================

def chunk_text_with_metadata(text: str, doc_name: str, doc_index: int,
                              size: int, overlap: int) -> List[Dict]:
    """Split text into chunks, tracking page boundaries from [Page X] markers."""
    # Split by [Page X] markers
    page_pattern = re.compile(r'\[Page\s+(\d+)\]')
    segments = page_pattern.split(text)

    # Build list of (page_number, page_text) pairs
    # segments alternates: text_before_first_marker, page_num, text, page_num, text, ...
    page_texts = []
    if segments and segments[0].strip():
        # Text before any page marker -> page 0 (unknown)
        page_texts.append((0, segments[0].strip()))

    i = 1
    while i < len(segments) - 1:
        page_num = int(segments[i])
        page_content = segments[i + 1].strip() if i + 1 < len(segments) else ""
        if page_content:
            page_texts.append((page_num, page_content))
        i += 2

    if not page_texts:
        # No page markers found (e.g., docx, xlsx) — treat as single page
        words = text.split()
        if not words or len(" ".join(words)) <= 20:
            return []
        raw_chunks = _chunk_words(words, size, overlap)
        return [{"text": c, "doc_name": doc_name, "pages": [], "doc_index": doc_index}
                for c in raw_chunks]

    # Chunk within and across pages, preserving page attribution
    chunks = []
    # Accumulate words with page tracking
    all_words = []  # list of (word, page_num)
    for page_num, page_content in page_texts:
        for w in page_content.split():
            all_words.append((w, page_num))

    if not all_words:
        return []

    if len(all_words) <= size:
        text_str = " ".join(w for w, _ in all_words)
        if len(text_str) > 20:
            pages = sorted(set(p for _, p in all_words if p > 0))
            chunks.append({"text": text_str, "doc_name": doc_name,
                           "pages": pages, "doc_index": doc_index})
        return chunks

    start = 0
    while start < len(all_words):
        end = min(start + size, len(all_words))
        # Try to break at sentence boundary
        if end < len(all_words):
            for j in range(end, max(end - size // 5, start + 1), -1):
                if all_words[j - 1][0][-1] in '.!?':
                    end = j
                    break
        chunk_words = all_words[start:end]
        text_str = " ".join(w for w, _ in chunk_words).strip()
        if len(text_str) > 20:
            pages = sorted(set(p for _, p in chunk_words if p > 0))
            chunks.append({"text": text_str, "doc_name": doc_name,
                           "pages": pages, "doc_index": doc_index})
        start += max(end - start - overlap, 1)

    return chunks


def _chunk_words(words: List[str], size: int, overlap: int) -> List[str]:
    """Simple word-based chunking without metadata (helper)."""
    if len(words) <= size:
        return [" ".join(words)] if len(" ".join(words)) > 20 else []
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + size, len(words))
        if end < len(words):
            for j in range(end, max(end - size // 5, start + 1), -1):
                if words[j - 1][-1] in '.!?': end = j; break
        c = " ".join(words[start:end]).strip()
        if len(c) > 20: chunks.append(c)
        start += max(end - start - overlap, 1)
    return chunks


def process_file(fp: Path, doc_index: int = 0) -> List[Dict]:
    """Process a file and return list of chunk dicts with metadata."""
    try:
        raw = _extract_text(fp)
        if not raw or len(raw.strip()) < 30: return []
        text = clean_text(raw)
        doc_name = fp.name  # Keep original filename for matching
        chunks = chunk_text_with_metadata(text, doc_name, doc_index,
                                           CHUNK_SIZE, CHUNK_OVERLAP)
        return chunks[:MAX_CHUNKS_PER_DOC]
    except Exception as e:
        logger.error(f"Err {fp.name}: {e}"); return []


# ==============================================================================
# TOKENIZATION + EXPANSION
# ==============================================================================

_STOPWORDS = frozenset(
    "the a an is are was were be been being have has had do does did will would could "
    "should may might shall can need must ought i me my we our you your he she it they "
    "them their this that these those which what who whom when where how why if then than "
    "but and or not no nor so too very just also as at by for from in into of on to up "
    "with about after before between during through above below each all any both few more "
    "most other some such only own same here there its over under out because until while "
    "again further once".split()
)
_TOKEN_RE = re.compile(r'[a-z0-9$%]+(?:[.\'-][a-z0-9]+)*')

def tokenize(text: str) -> List[str]:
    return [t for t in _TOKEN_RE.findall(text.lower()) if t not in _STOPWORDS and len(t) > 1]

_EXPANSIONS = {
    "revenue": {"sales","income","earnings","turnover","profit","financial"},
    "who": {"name","person","officer","director"},
    "when": {"date","year","month"},
    "where": {"location","headquarters","city","country"},
    "how much": {"amount","total","value","cost"},
    "how many": {"number","count","total"},
    "governing law": {"jurisdiction","governed","applicable","delaware","york"},
    "bench": {"justice","judge","opinion","court","delivered","justices"},
    "anticompetitive": {"competition","merger","combination","threshold","herfindahl","hhindex"},
    "notify": {"notification","filing","threshold","approval"},
    "scotus": {"supreme","court","case","opinion","cases"},
    "turnover": {"revenue","sales","threshold","crore","billion","assets"},
}

def expand_query(q: str) -> str:
    ql = q.lower()
    extras = set()
    for trigger, kws in _EXPANSIONS.items():
        if trigger in ql: extras.update(kws)
    return f"{q} {' '.join(extras)}" if extras else q


# ==============================================================================
# BM25
# ==============================================================================

class BM25Index:
    def __init__(self, k1=1.5, b=0.75):
        self.k1, self.b = k1, b
        self.doc_freqs, self.doc_lens = {}, []
        self.avg_dl, self.n_docs = 0.0, 0
        self.inverted = {}

    def build(self, docs: List[str]):
        self.n_docs = len(docs); total = 0
        for did, doc in enumerate(docs):
            toks = tokenize(doc)
            self.doc_lens.append(len(toks)); total += len(toks)
            tf = {}
            for t in toks: tf[t] = tf.get(t, 0) + 1
            for t, c in tf.items():
                self.doc_freqs[t] = self.doc_freqs.get(t, 0) + 1
                self.inverted.setdefault(t, []).append((did, c))
        self.avg_dl = total / max(self.n_docs, 1)

    def query(self, text, top_k=20):
        scores = {}
        for t in tokenize(text):
            if t not in self.inverted: continue
            df = self.doc_freqs.get(t, 0)
            idf = math.log((self.n_docs - df + 0.5) / (df + 0.5) + 1.0)
            for did, tf in self.inverted[t]:
                dl = self.doc_lens[did]
                tfn = (tf * (self.k1 + 1)) / (tf + self.k1 * (1 - self.b + self.b * dl / self.avg_dl))
                scores[did] = scores.get(did, 0.0) + idf * tfn
        return [(s, d) for d, s in sorted(scores.items(), key=lambda x: -x[1])[:top_k]]

    def get_idf(self, t):
        df = self.doc_freqs.get(t, 0)
        return math.log((self.n_docs - df + 0.5) / (df + 0.5) + 1.0) if df else 0.0

    def to_dict(self):
        return {"k1": self.k1, "b": self.b, "doc_freqs": self.doc_freqs,
                "doc_lens": self.doc_lens, "avg_dl": self.avg_dl, "n_docs": self.n_docs,
                "inverted": dict(self.inverted)}

    @classmethod
    def from_dict(cls, d):
        o = cls(d["k1"], d["b"])
        o.doc_freqs, o.doc_lens, o.avg_dl, o.n_docs = d["doc_freqs"], d["doc_lens"], d["avg_dl"], d["n_docs"]
        o.inverted = {k: [tuple(x) for x in v] for k, v in d["inverted"].items()}
        return o


class TFIDFVectorizer:
    def __init__(self, bm25, min_df=2):
        self.bm25 = bm25
        self.vocab = {t: i for i, t in enumerate(sorted(
            t for t, df in bm25.doc_freqs.items() if df >= min_df))}
        self.vocab_size = len(self.vocab)

    def vectorize_documents(self, docs: List[str]):
        m = np.zeros((len(docs), self.vocab_size), dtype=np.float32)
        for did, doc in enumerate(docs):
            for t, c in Counter(tokenize(doc)).items():
                if t in self.vocab: m[did, self.vocab[t]] = (1+math.log(c)) * self.bm25.get_idf(t)
        n = np.linalg.norm(m, axis=1, keepdims=True); n[n==0] = 1.0
        return m / n

    def vectorize_query(self, text):
        v = np.zeros(self.vocab_size, dtype=np.float32)
        for t, c in Counter(tokenize(text)).items():
            if t in self.vocab: v[self.vocab[t]] = (1+math.log(c)) * self.bm25.get_idf(t)
        n = np.linalg.norm(v)
        return v/n if n > 0 else v


# ==============================================================================
# RETRIEVAL
# ==============================================================================

def hybrid_retrieve(bm25, chunks: List[Dict], query, top_k=10, bm25_k=80,
                    tfidf_matrix=None, tfidf_vec=None, **kw):
    """Retrieve top chunks. chunks is list of dicts with 'text' key."""
    expanded = expand_query(query)
    results = bm25.query(expanded, top_k=bm25_k)
    if not results: results = bm25.query(query, top_k=bm25_k)
    if not results: return []

    cands = [d for _, d in results]
    bm25_sc = {d: s for s, d in results}
    mx = max(bm25_sc.values()) or 1.0

    tfidf_sc = {}
    if tfidf_matrix is not None and tfidf_vec is not None:
        qv = tfidf_vec.vectorize_query(expanded)
        sims = tfidf_matrix[cands] @ qv
        for j, d in enumerate(cands): tfidf_sc[d] = float(sims[j])

    q_proper = set()
    for w in query.split():
        c = re.sub(r'[^a-zA-Z]', '', w)
        if c and c[0].isupper() and len(c) > 2: q_proper.add(c.lower())
    q_nums = set(re.findall(r'\b\d{4}\b|\$[\d,.]+|\b\d+(?:\.\d+)?%?\b', query))
    ql_words = query.lower().split()

    scored = []
    for d in cands:
        cl = chunks[d]["text"].lower()
        s = 0.45 * (bm25_sc[d] / mx) + 0.25 * tfidf_sc.get(d, 0.0)
        if q_proper: s += min(sum(0.06 for p in q_proper if p in cl), 0.24)
        if q_nums: s += min(sum(0.05 for n in q_nums if n.lower() in cl), 0.15)
        for n in range(min(len(ql_words), 5), 2, -1):
            for i in range(len(ql_words) - n + 1):
                if " ".join(ql_words[i:i+n]) in cl: s += 0.03 * n
        scored.append((s, d))

    scored.sort(reverse=True, key=lambda x: x[0])
    seen, out = set(), []
    for _, d in scored:
        key = chunks[d]["text"][:150]
        if key not in seen: seen.add(key); out.append(d)
        if len(out) >= top_k: break
    return out


# ==============================================================================
# CORPUS SUMMARY
# ==============================================================================

def build_corpus_summary(chunks: List[Dict]) -> str:
    """Build a summary of all documents in the corpus for meta-questions."""
    doc_names = sorted(set(c["doc_name"] for c in chunks))
    lines = [f"Total documents: {len(doc_names)}"]
    for name in doc_names:
        lines.append(f"  - {name}")
    return "\n".join(lines)


# ==============================================================================
# CLI
# ==============================================================================

@click.group()
def cli(): pass


@cli.command()
def ingest():
    global START_TIME; START_TIME = time.time()
    logger.info("=== INGEST START ===")
    dp = Path(DATA_DIR)
    if not dp.exists(): logger.error(f"{DATA_DIR} not found"); sys.exit(1)

    files = [f for f in dp.iterdir()
             if f.is_file() and not f.name.startswith('.') and f.name != 'Testing Set Questions.xlsx']
    files = sorted(files, key=lambda f: f.name)
    logger.info(f"Found {len(files)} files")
    adaptive_optimization(len(files))

    all_chunks = []
    limit = 23
    with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_WORKERS) as pool:
        futs = {pool.submit(process_file, f, idx): f for idx, f in enumerate(files)}
        done = 0
        for fut in concurrent.futures.as_completed(futs):
            if time.time() - START_TIME > limit:
                logger.warning(f"Time limit at {done}/{len(files)}")
                for f in futs: f.cancel()
                break
            try: all_chunks.extend(fut.result(timeout=max(1, limit-(time.time()-START_TIME))))
            except: pass
            done += 1
            if done % 20 == 0: logger.info(f"  {done}/{len(files)}, {len(all_chunks)} chunks")

    if not all_chunks: logger.error("No chunks!"); sys.exit(1)
    logger.info(f"{len(all_chunks)} chunks extracted")

    # Extract text list for BM25/TF-IDF indexing
    chunk_texts = [c["text"] for c in all_chunks]

    bm25 = BM25Index(); bm25.build(chunk_texts)
    vec = TFIDFVectorizer(bm25, min_df=2)
    mat = vec.vectorize_documents(chunk_texts)
    logger.info(f"BM25: {bm25.n_docs} docs, {len(bm25.doc_freqs)} terms. TF-IDF: {mat.shape}")

    with open(CHUNKS_PATH, 'w') as f: json.dump(all_chunks, f, ensure_ascii=False)
    with open(TFIDF_PATH, 'w') as f: json.dump(bm25.to_dict(), f, ensure_ascii=False)
    np.savez_compressed(TFIDF_VECTORS_PATH, matrix=mat)
    with open(VOCAB_PATH, 'w') as f: json.dump(vec.vocab, f)
    with open(META_PATH, 'w') as f:
        json.dump({"ingest_time": time.time()-START_TIME, "doc_count": len(files),
                    "chunk_count": len(all_chunks), "vocab_size": vec.vocab_size}, f)

    logger.info(f"=== INGEST DONE: {len(all_chunks)} chunks, {time.time()-START_TIME:.2f}s ===")
    gc.collect()


@cli.command()
def run():
    global START_TIME; START_TIME = time.time()
    logger.info("=== RUN START ===")

    gemini_ok = test_gemini() if GEMINI_API_KEY else False

    with open(CHUNKS_PATH) as f: chunks = json.load(f)
    with open(TFIDF_PATH) as f: bm25 = BM25Index.from_dict(json.load(f))

    # Extract text list for TF-IDF queries
    chunk_texts = [c["text"] for c in chunks]

    tfidf_mat, tfidf_vec = None, None
    if os.path.exists(TFIDF_VECTORS_PATH) and os.path.exists(VOCAB_PATH):
        try:
            tfidf_mat = np.load(TFIDF_VECTORS_PATH)['matrix']
            with open(VOCAB_PATH) as f: vocab = json.load(f)
            tfidf_vec = TFIDFVectorizer.__new__(TFIDFVectorizer)
            tfidf_vec.bm25, tfidf_vec.vocab, tfidf_vec.vocab_size = bm25, vocab, len(vocab)
        except: pass

    ingest_time = 0
    if os.path.exists(META_PATH):
        with open(META_PATH) as f: ingest_time = json.load(f).get("ingest_time", 0)

    # Build corpus summary for meta-questions
    corpus_summary = build_corpus_summary(chunks)

    import pandas as pd
    qdf = pd.read_excel(os.path.join(DATA_DIR, 'Testing Set Questions.xlsx'))
    questions = qdf['Question'].tolist()
    gt_col = next((c for c in ['Answer','answer','Answers','Expected Answer'] if c in qdf.columns), None)
    ground_truth = qdf[gt_col].tolist() if gt_col else [None]*len(questions)

    # Check for "Documents Referred" column for ground truth doc/page info
    doc_ref_col = next((c for c in ['Documents Referred', 'documents referred',
                                     'Document Referred', 'Doc Referred'] if c in qdf.columns), None)
    doc_refs = qdf[doc_ref_col].tolist() if doc_ref_col else [None]*len(questions)

    logger.info(f"{len(questions)} questions")

    # Retrieve — BM25 operates on text, but we pass chunk dicts for scoring
    retrieval = []
    for q in questions:
        # BM25 query uses chunk_texts internally via inverted index (already indexed by position)
        idxs = hybrid_retrieve(bm25, chunks, q, top_k=10, bm25_k=80,
                               tfidf_matrix=tfidf_mat, tfidf_vec=tfidf_vec)
        ctx_chunks = [chunks[i] for i in idxs]
        prompt = build_qa_prompt(q, ctx_chunks, corpus_summary=corpus_summary)
        retrieval.append({"q": q, "ctx": ctx_chunks, "prompt": prompt})
    logger.info(f"Retrieval: {time.time()-START_TIME:.2f}s")

    # Gemini — all in parallel
    raw_responses = [None] * len(questions)
    if gemini_ok:
        def _call(i): return i, call_gemini(retrieval[i]["prompt"], timeout=20.0)
        with concurrent.futures.ThreadPoolExecutor(max_workers=min(10, len(questions))) as pool:
            futs = {pool.submit(_call, i): i for i in range(len(questions))}
            for fut in concurrent.futures.as_completed(futs):
                try:
                    i, a = fut.result(timeout=20.0)
                    if a: raw_responses[i] = a
                except: pass

    # Results
    final = []
    for i, q in enumerate(questions):
        raw_resp = raw_responses[i] or ""
        parsed = parse_gemini_response(raw_resp)
        ans = parsed["answer"] or "Information not found in provided documents"
        src = "gemini" if raw_responses[i] else "no_answer"

        # If Gemini didn't return source info, infer from top retrieved chunks
        doc_names = parsed["doc_names"]
        page_numbers = parsed["page_numbers"]
        if not doc_names and retrieval[i]["ctx"]:
            # Use the top chunk's metadata as fallback
            top_chunk = retrieval[i]["ctx"][0]
            doc_names = [top_chunk["doc_name"]]
            page_numbers = top_chunk.get("pages", [])

        gt = ground_truth[i] if i < len(ground_truth) else None
        gt_s = str(gt).strip() if gt and str(gt).strip().lower() not in ('nan','none','') else None
        sc = score_answer(ans, gt_s) if gt_s else None

        # Ground truth doc reference
        gt_doc_ref = doc_refs[i] if i < len(doc_refs) else None
        gt_doc_ref_s = str(gt_doc_ref).strip() if gt_doc_ref and str(gt_doc_ref).strip().lower() not in ('nan','none','') else None

        final.append({
            "question": q,
            "answer": ans,
            "document_name": doc_names,
            "page_numbers": page_numbers,
            "ground_truth": gt_s,
            "ground_truth_doc_ref": gt_doc_ref_s,
            "source": src,
            "scores": sc
        })
        logger.info(f"Q{i+1} [{src}]: {ans[:120]}")
        if doc_names:
            logger.info(f"  -> Docs: {doc_names}, Pages: {page_numbers}")

    with open(RESULTS_PATH, 'w') as f: json.dump(final, f, indent=2, ensure_ascii=False)
    logger.info(f"=== RUN DONE: {time.time()-START_TIME:.2f}s, total={ingest_time+time.time()-START_TIME:.2f}s ===")
    gc.collect()


if __name__ == "__main__":
    cli()
